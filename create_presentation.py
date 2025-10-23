from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.alignment = PP_ALIGN.CENTER
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(100, 100, 100)

    return slide

def add_content_slide(prs, title):
    """Add a slide with title and content area"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)

    return slide

# Slide 1: Title
add_title_slide(prs, "Hospital Pyrolysis System", "Mass Balance Model Overview")

# Slide 2: Plant Structure Overview
slide = add_content_slide(prs, "Overall Plant Structure")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

components = [
    ("1. Feed System", "Mulberry wood chips + water + treatment chemicals"),
    ("2. Pyrolyzer (300°C)", "Thermal decomposition into char, bio-gas, and bio-oil"),
    ("3. Char Separator", "Separates solid char from gas streams"),
    ("4. PBR Filter", "Activated char filter for grey water treatment"),
    ("5. Buffer Tank", "Grey water collection and distribution"),
    ("6. Inert Gas System", "Carrier gas for pyrolysis process")
]

for i, (comp, desc) in enumerate(components):
    p = tf.add_paragraph()
    p.text = comp
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = f"   {desc}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(12)

# Remove first empty paragraph
if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 3: Process Flow
slide = add_content_slide(prs, "Process Flow Diagram")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

flow_text = """
Feed Input [16]
    ↓
Pyrolyzer [13] ← Inert Gas [14]
    ↓
Char Separator
    ├─→ Char [5]
    │       ↓
    │   PBR Filter [4]
    │       ↓
    │   Grey Water Treatment [1,2,3]
    │
    └─→ Gas Stream [6]
"""

p = tf.add_paragraph()
p.text = flow_text
p.font.name = "Courier New"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 51, 102)

# Slide 4: Mass Balance Tool Architecture
slide = add_content_slide(prs, "Mass Balance Tool Architecture")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

sections = [
    ("MATLAB System Composer", [
        "Hospital_Context.slx - Visual system architecture",
        "unitProfile.xml - Property definitions with stereotypes",
        "M_* stereotypes track mass balance properties"
    ]),
    ("MATLAB Scripts", [
        "MassBalanceMaster.m - Complete workflow orchestrator",
        "VariableList.m - User input variables",
        "ExportToMaple.m - Export to Maple format",
        "LoadMapleValues.m - Import calculated values"
    ]),
    ("Maple Environment", [
        "MassBalanceVariables.mpl - Auto-generated variables",
        "MassBalanceCalculation.mw - User equations",
        "ExportMapleToMatlab.mpl - Export results to MATLAB"
    ]),
    ("Output", [
        "MapleExportedVariables.m - Calculated results",
        "Updated System Composer model with solved values"
    ])
]

for section, items in sections:
    p = tf.add_paragraph()
    p.text = section
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(6)

    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"  • {item}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(60, 60, 60)
        p2.space_after = Pt(3)

    p2.space_after = Pt(10)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 5: File Interaction Workflow
slide = add_content_slide(prs, "File Interaction Workflow")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

workflow = [
    ("Step 1: Initialize", "VariableList.m → Hospital_Context.slx"),
    ("Step 2: Clear Model", "Clear all M_* properties in connectors"),
    ("Step 3: Write Initial Values", "Write user variables into model"),
    ("Step 4: Export to Maple", "ExportToMaple.m → MassBalanceVariables.mpl"),
    ("Step 5: Calculate", "Maple solves equations in MassBalanceCalculation.mw"),
    ("Step 6: Export Results", "ExportMapleToMatlab.mpl → MapleExportedVariables.m"),
    ("Step 7: Import Results", "LoadMapleValues.m reads calculated values"),
    ("Step 8: Update Model", "Write calculated values back to Hospital_Context.slx")
]

for i, (step, desc) in enumerate(workflow, 1):
    p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = f"   {desc}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(10)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 6: Stream [16] - Feed Input
slide = add_content_slide(prs, "Stream [16]: Feed Input to Pyrolyzer")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Stream Properties:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(10)

stream_data = [
    ("Total Mass Flow Rate", "1.5385 kg/s"),
    ("Chip Mass Flow Rate", "1.0 kg/s (65%)"),
    ("Water (H2O)", "0.3077 kg/s (20%)"),
    ("Treatment Chemicals", "0.2308 kg/s (15%)"),
    ("", ""),
    ("Component Mass Fractions:", ""),
    ("  • Chips", "65%"),
    ("  • Water", "20%"),
    ("  • Treatment Chemicals", "15%")
]

for label, value in stream_data:
    p = tf.add_paragraph()
    if label == "":
        p.space_after = Pt(5)
        continue
    if "Component" in label:
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
    elif label.startswith("  "):
        p.text = f"{label}: {value}"
        p.font.size = Pt(14)
        p.level = 1
    else:
        p.text = f"{label}: {value}"
        p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(6)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 7: Stream [14] - Inert Gas Input
slide = add_content_slide(prs, "Stream [14]: Inert Gas Supply")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Stream Properties:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(10)

stream_data = [
    ("Total Mass Flow Rate", "0.5 kg/s"),
    ("Inert Mass Flow Rate", "0.5 kg/s (100%)"),
    ("", ""),
    ("Purpose:", "Carrier gas for pyrolysis process"),
    ("", "Provides inert atmosphere for thermal decomposition")
]

for label, value in stream_data:
    p = tf.add_paragraph()
    if label == "":
        p.space_after = Pt(5)
        continue
    if ":" in label and "Purpose" in label:
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
    else:
        p.text = f"{label}: {value}" if label and not label.startswith("  ") else value
        p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(6)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 8: Stream [13] - Pyrolyzer Output
slide = add_content_slide(prs, "Stream [13]: Pyrolyzer Output")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Operating Conditions:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "Temperature: 300°C  |  Density: 200 kg/m³"
p.font.size = Pt(14)
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "Stream Composition:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(8)

stream_data = [
    ("Total Mass Flow Rate", "2.0385 kg/s", ""),
    ("", "", ""),
    ("Char", "0.45 kg/s", "22.08%"),
    ("Bio-gas (Gas1)", "0.25 kg/s", "12.26%"),
    ("Bio-gas (Gas2)", "0.20 kg/s", "9.81%"),
    ("Water (H2O)", "0.3077 kg/s", "15.09%"),
    ("Treatment Chemicals", "0.2308 kg/s", "11.32%"),
    ("Inert Gas", "0.5 kg/s", "24.53%"),
    ("", "", ""),
    ("Conversion Fractions (from chips):", "", ""),
    ("  • Char", "45%", ""),
    ("  • Gas1", "25%", ""),
    ("  • Gas2", "20%", "")
]

for label, value, pct in stream_data:
    p = tf.add_paragraph()
    if label == "":
        p.space_after = Pt(5)
        continue
    if "Conversion" in label:
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
    elif label.startswith("  "):
        p.text = f"{label}: {value}"
        p.font.size = Pt(13)
        p.level = 1
    else:
        display_text = f"{label}: {value}"
        if pct:
            display_text += f" ({pct})"
        p.text = display_text
        p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(4)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 9: Stream [5] - Pure Char Output
slide = add_content_slide(prs, "Stream [5]: Char Separator - Char Output")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Stream Properties:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(10)

stream_data = [
    ("Total Mass Flow Rate", "0.4500 kg/s"),
    ("Char Mass Flow Rate", "0.45 kg/s"),
    ("Water (H2O)", "4.5×10⁻⁹ kg/s (trace)"),
    ("", ""),
    ("Purity:", ""),
    ("  • Char Mass Fraction", "~100% (0.99999999)"),
    ("  • Water Mass Fraction", "~0% (1×10⁻⁸)"),
    ("", ""),
    ("Notes:", ""),
    ("  Char is extremely dry, suitable for use in PBR filter"),
    ("  Water content minimized to near-zero through separation")
]

for item in stream_data:
    p = tf.add_paragraph()
    if isinstance(item, str):
        label = item
        value = ""
    else:
        label, value = item if len(item) == 2 else (item[0], item[1])

    if label == "":
        p.space_after = Pt(5)
        continue
    if ":" in label and not label.startswith("  "):
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
    elif label.startswith("  •"):
        p.text = f"{label}: {value}"
        p.font.size = Pt(14)
        p.level = 1
    elif label.startswith("  "):
        p.text = label
        p.font.size = Pt(13)
        p.level = 1
        p.font.italic = True
    else:
        p.text = f"{label}: {value}" if value else label
        p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(6)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 10: Stream [6] - Gas Stream Output
slide = add_content_slide(prs, "Stream [6]: Char Separator - Gas Output")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Operating Conditions:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "Temperature: 290°C  |  Density: 180 kg/m³"
p.font.size = Pt(14)
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "Stream Composition:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(8)

stream_data = [
    ("Total Mass Flow Rate", "1.5885 kg/s", ""),
    ("Volume Flow Rate", "0.008825 m³/s", ""),
    ("", "", ""),
    ("Bio-gas (Gas1)", "0.25 kg/s", "15.74%"),
    ("Bio-gas (Gas2)", "0.20 kg/s", "12.59%"),
    ("Water (H2O)", "0.3077 kg/s", "19.37%"),
    ("Treatment Chemicals", "0.2308 kg/s", "14.53%"),
    ("Inert Gas", "0.5 kg/s", "31.48%")
]

for label, value, pct in stream_data:
    p = tf.add_paragraph()
    if label == "":
        p.space_after = Pt(5)
        continue
    display_text = f"{label}: {value}"
    if pct:
        display_text += f" ({pct})"
    p.text = display_text
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(6)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 11: Stream [4] - PBR Filter Output
slide = add_content_slide(prs, "Stream [4]: PBR Filter Output")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Stream Properties:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(10)

stream_data = [
    ("Total Mass Flow Rate", "0.4683 kg/s"),
    ("", ""),
    ("Component Breakdown:", ""),
    ("  • Char", "0.45 kg/s (96.1%)"),
    ("  • Water (H2O)", "0.01405 kg/s (3.0%)"),
    ("  • CHx (Contaminants)", "0.004214 kg/s (0.9%)"),
    ("", ""),
    ("Mass Fractions:", ""),
    ("  • Char", "0.961"),
    ("  • Water", "0.03"),
    ("  • CHx", "0.009"),
    ("", ""),
    ("Function:", ""),
    ("  Spent char filter containing adsorbed CHx contaminants"),
    ("  and absorbed water from grey water treatment process")
]

for item in stream_data:
    p = tf.add_paragraph()
    if isinstance(item, str):
        label = item
        value = ""
    else:
        label, value = item if len(item) == 2 else (item[0], item[1])

    if label == "":
        p.space_after = Pt(5)
        continue
    if ":" in label and not label.startswith("  "):
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
    elif label.startswith("  •"):
        p.text = f"{label}: {value}" if value else label
        p.font.size = Pt(14)
        p.level = 1
    elif label.startswith("  "):
        p.text = label
        p.font.size = Pt(13)
        p.level = 1
        p.font.italic = True
    else:
        p.text = f"{label}: {value}" if value else label
        p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(6)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 12: Streams [1,2,3] - Grey Water System
slide = add_content_slide(prs, "Streams [1,2,3]: Grey Water Treatment System")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Operating Conditions: 10°C, Density: 999.7 kg/m³"
p.font.size = Pt(14)
p.font.bold = True
p.space_after = Pt(12)

streams = [
    ("Stream [1]: Buffer Tank Input", [
        ("Mass Flow Rate", "0.9997 kg/s"),
        ("Volume Flow Rate", "1.0 m³/s"),
        ("CHx Concentration", "1.0125 mg/L"),
        ("Temperature", "10°C"),
        ("Density", "999.7 kg/m³")
    ]),
    ("Stream [2]: To PBR Filter", [
        ("Mass Flow Rate", "0.9997 kg/s"),
        ("Volume Flow Rate", "1.0 m³/s"),
        ("CHx Concentration", "1.0125 mg/L"),
        ("(Same as Stream [1])", "")
    ]),
    ("Stream [3]: Filter Outlet (Treated Water)", [
        ("Mass Flow Rate", "0.9814 kg/s"),
        ("Volume Flow Rate", "0.9859 m³/s"),
        ("CHx Concentration", "0.01 mg/L (99% reduction!)"),
        ("Water (H2O) Flow", "0.9857 kg/s"),
        ("Fluid Density", "999.7 kg/m³")
    ])
]

for stream_name, data in streams:
    p = tf.add_paragraph()
    p.text = stream_name
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(6)

    for label, value in data:
        if value:
            p = tf.add_paragraph()
            p.text = f"  {label}: {value}"
            p.font.size = Pt(12)
            p.level = 1
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.space_after = Pt(4)

    p.space_after = Pt(10)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 13: Mass Balance Summary
slide = add_content_slide(prs, "Overall Mass Balance Summary")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "System Inputs:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(8)

inputs = [
    ("Feed Stream [16]", "1.5385 kg/s"),
    ("Inert Gas [14]", "0.5 kg/s"),
    ("Grey Water [1]", "0.9997 kg/s"),
    ("TOTAL INPUT", "3.0382 kg/s")
]

for label, value in inputs:
    p = tf.add_paragraph()
    if "TOTAL" in label:
        p.text = f"{label}: {value}"
        p.font.bold = True
        p.font.size = Pt(14)
    else:
        p.text = f"  {label}: {value}"
        p.font.size = Pt(13)
        p.level = 1
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(5)

p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "System Outputs:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.space_after = Pt(8)

outputs = [
    ("Gas Stream [6]", "1.5885 kg/s"),
    ("Spent Char [4]", "0.4683 kg/s"),
    ("Treated Water [3]", "0.9814 kg/s"),
    ("TOTAL OUTPUT", "3.0382 kg/s")
]

for label, value in outputs:
    p = tf.add_paragraph()
    if "TOTAL" in label:
        p.text = f"{label}: {value}"
        p.font.bold = True
        p.font.size = Pt(14)
    else:
        p.text = f"  {label}: {value}"
        p.font.size = Pt(13)
        p.level = 1
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.space_after = Pt(5)

p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "✓ Mass Balance Verified: Input = Output"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Slide 14: Key Performance Indicators
slide = add_content_slide(prs, "Key Performance Indicators")
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

kpis = [
    ("Pyrolysis Efficiency", [
        "Char Yield: 45% (from wood chips)",
        "Bio-gas Yield: 45% (Gas1 + Gas2)",
        "Operating Temperature: 300°C"
    ]),
    ("Char Separation Performance", [
        "Char Purity: 99.99999%",
        "Moisture Content: <0.000001%",
        "Recovery Rate: 100%"
    ]),
    ("Grey Water Treatment", [
        "CHx Removal: 99.01% (1.0125 → 0.01 mg/L)",
        "Water Recovery: 98.17%",
        "Char Adsorption Capacity: 0.9% CHx"
    ]),
    ("Overall System", [
        "Total Throughput: 3.04 kg/s",
        "Char Production: 0.45 kg/s",
        "Gas Production: 1.59 kg/s"
    ])
]

for category, metrics in kpis:
    p = tf.add_paragraph()
    p.text = category
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(6)

    for metric in metrics:
        p = tf.add_paragraph()
        p.text = f"  • {metric}"
        p.font.size = Pt(13)
        p.level = 1
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(4)

    p.space_after = Pt(10)

if tf.paragraphs[0].text == "":
    p_to_remove = tf.paragraphs[0]._element
    p_to_remove.getparent().remove(p_to_remove)

# Save presentation
output_path = r"C:\Users\jakeo\Documents\MATLAB\3rd_Year_Project_C3\Hospital_Pyrolysis_System_Presentation.pptx"
prs.save(output_path)

print(f"Presentation created successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")

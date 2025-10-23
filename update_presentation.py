"""
Update PowerPoint Presentation with New Values
This script reads new values from MapleExportedVariables.m and updates
the presentation while preserving all formatting.
"""

from pptx import Presentation
import re
from pathlib import Path

def parse_maple_variables(file_path):
    """Parse MapleExportedVariables.m and return a dictionary of variables"""
    variables = {}

    with open(file_path, 'r') as f:
        for line in f:
            # Match lines like: H2OMassFlowRate__(3) = 0.9856521332;
            match = re.match(r'(\w+)__\((\d+)\)\s*=\s*([\d.eE+-]+);', line)
            if match:
                var_name = match.group(1)
                connector = match.group(2)
                value = float(match.group(3))

                # Store with both formats for matching
                key = f"{var_name}__({connector})"
                key_bracket = f"{var_name}__[{connector}]"

                variables[key] = value
                variables[key_bracket] = value

    return variables

def format_value(value):
    """Format a numeric value for display"""
    # Handle very small numbers (scientific notation)
    if abs(value) < 0.00001 and value != 0:
        # Format in scientific notation
        return f"{value:.2e}"
    # Handle regular numbers
    elif abs(value) >= 1000:
        return f"{value:.0f}"
    elif abs(value) >= 100:
        return f"{value:.1f}"
    elif abs(value) >= 1:
        return f"{value:.4f}".rstrip('0').rstrip('.')
    else:
        return f"{value:.10f}".rstrip('0').rstrip('.')

def create_value_mapping():
    """Create a mapping of common derived values shown in the presentation"""

    # Parse variables
    maple_vars = parse_maple_variables(
        r"C:\Users\jakeo\Documents\MATLAB\3rd_Year_Project_C3\Model\BackendMassBalanceMaster\MapleExportedVariables.m"
    )

    # Direct mappings from variable names to display values
    mappings = {}

    # Stream [16] values
    if 'MassFlowRate__(16)' in maple_vars:
        mappings['1.5385'] = format_value(maple_vars['MassFlowRate__(16)'])
        mappings['1.538461538'] = format_value(maple_vars['MassFlowRate__(16)'])

    if 'ChipMassFlowRate__(16)' in maple_vars:
        mappings['1.0'] = format_value(maple_vars['ChipMassFlowRate__(16)'])

    if 'H2OMassFlowRate__(16)' in maple_vars:
        mappings['0.3077'] = format_value(maple_vars['H2OMassFlowRate__(16)'])
        mappings['0.3076923076'] = format_value(maple_vars['H2OMassFlowRate__(16)'])

    if 'TreatmentChemicalsMassFlowRate__(16)' in maple_vars:
        mappings['0.2308'] = format_value(maple_vars['TreatmentChemicalsMassFlowRate__(16)'])
        mappings['0.2307692307'] = format_value(maple_vars['TreatmentChemicalsMassFlowRate__(16)'])

    # Stream [14] values
    if 'InertMassFlowRate__(14)' in maple_vars:
        mappings['0.5'] = format_value(maple_vars['InertMassFlowRate__(14)'])

    if 'MassFlowRate__(14)' in maple_vars:
        mappings['0.5'] = format_value(maple_vars['MassFlowRate__(14)'])

    # Stream [13] values
    if 'MassFlowRate__(13)' in maple_vars:
        mappings['2.0385'] = format_value(maple_vars['MassFlowRate__(13)'])
        mappings['2.038461538'] = format_value(maple_vars['MassFlowRate__(13)'])

    if 'CharMassFlowRate__(13)' in maple_vars:
        mappings['0.45'] = format_value(maple_vars['CharMassFlowRate__(13)'])

    if 'Gas1MassFlowRate__(13)' in maple_vars:
        mappings['0.25'] = format_value(maple_vars['Gas1MassFlowRate__(13)'])

    if 'Gas2MassFlowRate__(13)' in maple_vars:
        mappings['0.20'] = format_value(maple_vars['Gas2MassFlowRate__(13)'])
        mappings['0.2'] = format_value(maple_vars['Gas2MassFlowRate__(13)'])

    # Stream [5] values
    if 'MassFlowRate__(5)' in maple_vars:
        mappings['0.4500'] = format_value(maple_vars['MassFlowRate__(5)'])
        mappings['0.450001'] = format_value(maple_vars['MassFlowRate__(5)'])

    if 'H2OMassFlowRate__(5)' in maple_vars:
        mappings['4.5×10⁻⁹'] = format_value(maple_vars['H2OMassFlowRate__(5)'])
        mappings['4.50001e-09'] = format_value(maple_vars['H2OMassFlowRate__(5)'])

    # Stream [6] values
    if 'MassFlowRate__(6)' in maple_vars:
        mappings['1.5885'] = format_value(maple_vars['MassFlowRate__(6)'])
        mappings['1.588460538'] = format_value(maple_vars['MassFlowRate__(6)'])

    if 'VolumeFlowRate__(6)' in maple_vars:
        mappings['0.008825'] = format_value(maple_vars['VolumeFlowRate__(6)'])
        mappings['0.008824780767'] = format_value(maple_vars['VolumeFlowRate__(6)'])

    # Stream [4] values
    if 'MassFlowRate__(4)' in maple_vars:
        mappings['0.4683'] = format_value(maple_vars['MassFlowRate__(4)'])
        mappings['0.4682622268'] = format_value(maple_vars['MassFlowRate__(4)'])

    if 'H2OMassFlowRate__(4)' in maple_vars:
        mappings['0.01405'] = format_value(maple_vars['H2OMassFlowRate__(4)'])
        mappings['0.0140478668'] = format_value(maple_vars['H2OMassFlowRate__(4)'])

    if 'CHxMassFlowRate__(4)' in maple_vars:
        mappings['0.004214'] = format_value(maple_vars['CHxMassFlowRate__(4)'])
        mappings['0.004214360041'] = format_value(maple_vars['CHxMassFlowRate__(4)'])

    # Streams [1,2,3] values
    if 'MassFlowRate__(1)' in maple_vars:
        mappings['0.9997'] = format_value(maple_vars['MassFlowRate__(1)'])

    if 'VolumeFlowRate__(1)' in maple_vars:
        mappings['1.0'] = format_value(maple_vars['VolumeFlowRate__(1)'])

    if 'CHxConcentration__(1)' in maple_vars:
        mappings['1.0125'] = format_value(maple_vars['CHxConcentration__(1)'])

    if 'MassFlowRate__(3)' in maple_vars:
        mappings['0.9814'] = format_value(maple_vars['MassFlowRate__(3)'])
        mappings['0.9814387732'] = format_value(maple_vars['MassFlowRate__(3)'])

    if 'VolumeFlowRate__(3)' in maple_vars:
        mappings['0.9859'] = format_value(maple_vars['VolumeFlowRate__(3)'])
        mappings['0.9859479176'] = format_value(maple_vars['VolumeFlowRate__(3)'])

    if 'CHxConcentration__(3)' in maple_vars:
        mappings['0.01'] = format_value(maple_vars['CHxConcentration__(3)'])

    if 'H2OMassFlowRate__(3)' in maple_vars:
        mappings['0.9857'] = format_value(maple_vars['H2OMassFlowRate__(3)'])
        mappings['0.9856521332'] = format_value(maple_vars['H2OMassFlowRate__(3)'])

    # Mass balance totals
    total_input = 0
    total_output = 0

    if all(k in maple_vars for k in ['MassFlowRate__(16)', 'MassFlowRate__(14)', 'MassFlowRate__(1)']):
        total_input = maple_vars['MassFlowRate__(16)'] + maple_vars['MassFlowRate__(14)'] + maple_vars['MassFlowRate__(1)']
        mappings['3.0382'] = format_value(total_input)

    if all(k in maple_vars for k in ['MassFlowRate__(6)', 'MassFlowRate__(4)', 'MassFlowRate__(3)']):
        total_output = maple_vars['MassFlowRate__(6)'] + maple_vars['MassFlowRate__(4)'] + maple_vars['MassFlowRate__(3)']
        mappings['3.0382'] = format_value(total_output)

    # Percentage values from mass fractions
    if 'MassFraction_Chips__(16)' in maple_vars:
        pct = maple_vars['MassFraction_Chips__(16)'] * 100
        mappings['65'] = f"{pct:.0f}" if pct == int(pct) else f"{pct:.1f}"

    if 'MassFraction_H2O__(16)' in maple_vars:
        pct = maple_vars['MassFraction_H2O__(16)'] * 100
        mappings['20'] = f"{pct:.0f}" if pct == int(pct) else f"{pct:.1f}"

    if 'MassFraction_TreatmentChemicals__(16)' in maple_vars:
        pct = maple_vars['MassFraction_TreatmentChemicals__(16)'] * 100
        mappings['15'] = f"{pct:.0f}" if pct == int(pct) else f"{pct:.1f}"

    # Conversion fractions
    if 'CharConversionFraction__(13)' in maple_vars:
        pct = maple_vars['CharConversionFraction__(13)'] * 100
        mappings['45'] = f"{pct:.0f}" if pct == int(pct) else f"{pct:.1f}"

    if 'Gas1ConversionFraction__(13)' in maple_vars:
        pct = maple_vars['Gas1ConversionFraction__(13)'] * 100
        mappings['25'] = f"{pct:.0f}" if pct == int(pct) else f"{pct:.1f}"

    if 'Gas2ConversionFraction__(13)' in maple_vars:
        pct = maple_vars['Gas2ConversionFraction__(13)'] * 100
        mappings['20'] = f"{pct:.0f}" if pct == int(pct) else f"{pct:.1f}"

    return mappings, maple_vars

def update_text_in_shape(shape, old_to_new_mapping):
    """Update text in a shape while preserving formatting"""
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            updated_text = original_text

            # Replace values using the mapping
            for old_val, new_val in old_to_new_mapping.items():
                # Use word boundaries to avoid partial replacements
                # But be flexible with context
                updated_text = re.sub(
                    r'\b' + re.escape(str(old_val)) + r'\b',
                    str(new_val),
                    updated_text
                )

            # Update the run text if it changed
            if updated_text != original_text:
                run.text = updated_text

def update_presentation(input_pptx, output_pptx=None):
    """
    Update presentation with new values from MapleExportedVariables.m

    Args:
        input_pptx: Path to the formatted presentation
        output_pptx: Path for output (if None, overwrites input)
    """

    # Load the presentation
    prs = Presentation(input_pptx)

    # Get value mappings
    mappings, maple_vars = create_value_mapping()

    print(f"Loaded presentation: {input_pptx}")
    print(f"Found {len(maple_vars)} variables in MapleExportedVariables.m")
    print(f"Created {len(mappings)} value mappings")
    print("\nUpdating slides...")

    # Iterate through all slides
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            update_text_in_shape(shape, mappings)

            # Handle grouped shapes
            if shape.shape_type == 6:  # Group
                for sub_shape in shape.shapes:
                    update_text_in_shape(sub_shape, mappings)

            # Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        update_text_in_shape(cell, mappings)

        print(f"  Updated slide {slide_num}")

    # Save the presentation
    output_path = output_pptx if output_pptx else input_pptx
    prs.save(output_path)

    print(f"\n[SUCCESS] Presentation updated successfully!")
    print(f"  Output: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path

if __name__ == "__main__":
    # Define paths
    presentation_path = r"C:\Users\jakeo\Documents\MATLAB\3rd_Year_Project_C3\Hospital_Pyrolysis_System_Presentation.pptx"

    # Option 1: Overwrite the existing presentation
    # update_presentation(presentation_path)

    # Option 2: Create a new updated version
    output_path = r"C:\Users\jakeo\Documents\MATLAB\3rd_Year_Project_C3\Hospital_Pyrolysis_System_Presentation_Updated.pptx"
    update_presentation(presentation_path, output_path)

    print("\nNote: To overwrite the original presentation instead of creating a new file,")
    print("      edit the script and use: update_presentation(presentation_path)")
